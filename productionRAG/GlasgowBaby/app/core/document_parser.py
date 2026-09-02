# app/core/document_parser.py

from __future__ import annotations
from logger import Logger

import csv
import io
import json
import logging
import re
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET

from pathlib import Path
from typing import Callable, Optional

import fitz  # PyMuPDF
import nltk
import pytesseract

from PIL import Image

# Optional dependencies used by specific formats.
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from striprtf.striprtf import rtf_to_text
except ImportError:
    rtf_to_text = None


# ============================================================================
# Configuration
# ============================================================================

DEFAULT_CHUNK_SIZE = 500
DEFAULT_CHUNK_OVERLAP = 100

# Images smaller than this are generally not useful as standalone OCR input.
MIN_IMAGE_WIDTH = 32
MIN_IMAGE_HEIGHT = 32

SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf",
    ".docx",
    ".doc",
    ".txt",
    ".text",
    ".md",
    ".markdown",
    ".html",
    ".htm",
    ".rtf",
    ".tex",
    ".latex",

    # Presentations
    ".pptx",
    ".ppt",

    # Spreadsheets / tabular
    ".xlsx",
    ".xls",
    ".csv",
    ".tsv",

    # Structured data
    ".json",
    ".xml",

    # Logs
    ".log",

    # Images
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".bmp",
    ".tiff",
    ".tif",
    ".gif",
}

# ============================================================================
# NLTK setup
# ============================================================================

def _ensure_nltk_resources() -> None:
    """
    Ensure the resources required for sentence tokenization exist.

    Modern NLTK versions require punkt_tab. Older installations may use
    punkt. We attempt both without making parser construction fail merely
    because the resource has not been downloaded yet.
    """

    resources = (
        ("tokenizers/punkt_tab/english/", "punkt_tab"),
        ("tokenizers/punkt/english.pickle", "punkt"),
    )

    for resource_path, package_name in resources:
        try:
            nltk.data.find(resource_path)
            return
        except LookupError:
            continue

    # Try the modern resource first.
    try:
        nltk.download(
            "punkt_tab",
            quiet=True,
            raise_on_error=False,
        )
        nltk.data.find("tokenizers/punkt_tab/english/")
        return
    except LookupError:
        pass

    # Fallback for older NLTK environments.
    try:
        nltk.download(
            "punkt",
            quiet=True,
            raise_on_error=False,
        )
        nltk.data.find("tokenizers/punkt/english.pickle")
        return
    except LookupError:
        pass

    raise RuntimeError(
        "NLTK sentence tokenizer resources are unavailable. "
        "Install/download 'punkt_tab' (or 'punkt' for older NLTK versions)."
    )


# ============================================================================
# Chunking
# ============================================================================

def chunk_text(
    text: str,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[str]:
    """
    Split text into semantically sensible chunks using NLTK sentence
    segmentation.

    Parameters
    ----------
    text:
        Text to chunk.

    chunk_size:
        Approximate maximum number of characters per chunk.

    overlap:
        Approximate number of characters from the previous chunk to carry
        into the next chunk.

    Returns
    -------
    list[str]
        A list of text chunks.

    Notes
    -----
    This deliberately preserves the same public return format as your
    original implementation: list[str].

    The chunker tries to:
      1. Preserve sentence boundaries.
      2. Keep chunks around chunk_size.
      3. Preserve paragraph boundaries where possible.
      4. Add overlap without splitting arbitrary sentences.
      5. Handle extremely long sentences safely.
    """

    if not isinstance(text, str):
        raise TypeError("text must be a string")

    if chunk_size <= 0:
        raise ValueError("chunk_size must be greater than 0")

    if overlap < 0:
        raise ValueError("overlap cannot be negative")

    if overlap >= chunk_size:
        raise ValueError(
            "overlap must be smaller than chunk_size"
        )

    text = _normalise_text(text)

    if not text:
        return []

    _ensure_nltk_resources()

    # Split paragraphs first so we retain some document structure.
    paragraphs = re.split(r"\n\s*\n+", text)

    sentences: list[str] = []

    for paragraph in paragraphs:
        paragraph = paragraph.strip()

        if not paragraph:
            continue

        try:
            paragraph_sentences = nltk.sent_tokenize(paragraph)
        except Exception:
            # A malformed/odd document should not necessarily kill ingestion.
            paragraph_sentences = [paragraph]

        for sentence in paragraph_sentences:
            sentence = sentence.strip()

            if sentence:
                sentences.append(sentence)

    if not sentences:
        return []

    chunks: list[str] = []
    current_sentences: list[str] = []
    current_length = 0

    for sentence in sentences:
        sentence_length = len(sentence)

        # A sentence larger than the target must be handled separately.
        if sentence_length > chunk_size:
            if current_sentences:
                chunks.append(" ".join(current_sentences).strip())
                current_sentences = []
                current_length = 0

            # Hard split an exceptionally long sentence.
            long_parts = _hard_split(
                sentence,
                chunk_size=chunk_size,
            )

            chunks.extend(long_parts)
            continue

        separator_length = 1 if current_sentences else 0
        proposed_length = (
            current_length
            + separator_length
            + sentence_length
        )

        if proposed_length <= chunk_size:
            current_sentences.append(sentence)
            current_length = proposed_length
            continue

        # Current chunk is full.
        if current_sentences:
            chunks.append(
                " ".join(current_sentences).strip()
            )

        # Build overlap from complete sentences where possible.
        overlap_sentences: list[str] = []
        overlap_length = 0

        if overlap > 0:
            for previous in reversed(current_sentences):
                addition = (
                    len(previous)
                    if not overlap_sentences
                    else len(previous) + 1
                )

                if overlap_length + addition > overlap:
                    break

                overlap_sentences.insert(0, previous)
                overlap_length += addition

        current_sentences = overlap_sentences + [sentence]

        current_length = len(
            " ".join(current_sentences)
        )

    if current_sentences:
        chunks.append(
            " ".join(current_sentences).strip()
        )

    # Final cleanup.
    return [
        chunk.strip()
        for chunk in chunks
        if chunk.strip()
    ]


def _hard_split(
    text: str,
    chunk_size: int,
) -> list[str]:
    """
    Safely split a sentence which is larger than chunk_size.

    Prefer whitespace boundaries before falling back to a hard character
    boundary.
    """

    parts: list[str] = []
    remaining = text.strip()

    while len(remaining) > chunk_size:
        split_at = remaining.rfind(
            " ",
            0,
            chunk_size + 1,
        )

        if split_at <= 0:
            split_at = chunk_size

        part = remaining[:split_at].strip()

        if part:
            parts.append(part)

        remaining = remaining[split_at:].strip()

    if remaining:
        parts.append(remaining)

    return parts


def _normalise_text(text: str) -> str:
    """
    Normalise text without destroying meaningful line structure.
    """

    text = text.replace("\x00", "")

    # Normalise different newline styles.
    text = text.replace("\r\n", "\n")
    text = text.replace("\r", "\n")

    # Collapse excessive horizontal whitespace.
    text = re.sub(r"[ \t]+", " ", text)

    # Avoid pathological blank-line runs.
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


# ============================================================================
# Document parser
# ============================================================================

class DocumentParser:
    """
    Production-oriented document extraction facade.

    Public usage:

        chunks = DocumentParser.process("document.pdf")

    or:

        parser = DocumentParser(logger=my_logger)
        chunks = parser.process("document.pdf")

    The parser always returns:

        list[str]

    Extraction errors are raised rather than silently returning corrupt
    content.
    """

    def __init__(
        self,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        overlap: int = DEFAULT_CHUNK_OVERLAP,
        ocr_enabled: bool = True,
    ) -> None:
        self.logger = Logger

        self.chunk_size = chunk_size
        self.overlap = overlap
        self.ocr_enabled = ocr_enabled

        # Explicit registry.
        #
        # This is what allows process() to determine whether a file type
        # is supported before attempting extraction.
        self.parsers: dict[str, Callable[[Path], str]] = {
            # Documents
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".doc": self._process_doc,

            # Text
            ".txt": self._process_text,
            ".text": self._process_text,
            ".md": self._process_markdown,
            ".markdown": self._process_markdown,
            ".log": self._process_text,

            # Web / markup
            ".html": self._process_html,
            ".htm": self._process_html,
            ".xml": self._process_xml,

            # Rich / typesetting
            ".rtf": self._process_rtf,
            ".tex": self._process_latex,
            ".latex": self._process_latex,

            # Presentations
            ".pptx": self._process_pptx,
            ".ppt": self._process_ppt,

            # Spreadsheets
            ".xlsx": self._process_xlsx,
            ".xls": self._process_xls,
            ".csv": self._process_csv,
            ".tsv": self._process_tsv,

            # Structured data
            ".json": self._process_json,

            # Images
            ".png": self._process_image,
            ".jpg": self._process_image,
            ".jpeg": self._process_image,
            ".webp": self._process_image,
            ".bmp": self._process_image,
            ".tiff": self._process_image,
            ".tif": self._process_image,
            ".gif": self._process_image,
        }

        self.logger.log("Initialised Document Parser")

    # ------------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------------

    def process(self, file_path: str | Path) -> list[str]:
        """
        Process a file and return only its chunks.

        Raises
        ------
        ValueError
            Unsupported file type.

        FileNotFoundError
            Input file does not exist.

        RuntimeError
            Extraction/conversion failure.
        """

        path = Path(file_path)

        if not path.exists():
            err = FileNotFoundError(
                f"File does not exist: {path}"
            )
            self.logger.exception(f"File does not exist: {path}", err)

            raise err

        if not path.is_file():
            err = ValueError(
                f"Path is not a file: {path}"
            )
            self.logger.exception(f"Path is not a file: {path}", err)
           
            raise err

        extension = path.suffix.lower()

        # --------------------------------------------------------------
        # This is the requested explicit supported-file check.
        # --------------------------------------------------------------
        if extension not in self.parsers:
            err = ValueError(
                f"Unsupported file type: {extension or '<none>'}"
            )
            self.logger.exception(f"Unsupported file type: {extension or '<none>'}", err)
                       
            raise err 

        parser = self.parsers[extension]

        self.logger.log(
            f"Processing {path.name} ({extension})"
        )

        try:
            extracted_text = parser(path)

            extracted_text = _normalise_text(
                extracted_text
            )

            if not extracted_text:
                self.logger.warn(
                    f"No text extracted from {path.name}"
                )
                return []

            chunks = chunk_text(
                extracted_text,
                chunk_size=self.chunk_size,
                overlap=self.overlap,
            )

            self.logger.log(
                f"Processed {path.name}: "
                f"{len(chunks)} chunks"
            )

            return chunks

        except Exception as exc:
            self.logger.exception(
                f"Failed to process {path.name}",
                exc,
            )
            raise RuntimeError(
                f"Failed to process {path}: "
                f"{type(exc).__name__}: {exc}"
            ) from exc

    # ------------------------------------------------------------------------
    # PDF
    # ------------------------------------------------------------------------

    def _process_pdf(self, path: Path) -> str:
        """
        Extract native PDF text and OCR pages that contain little/no text.

        Embedded images are therefore covered indirectly: if a PDF page is
        essentially an image/scanned page, OCR is performed on the rendered
        page.
        """

        output: list[str] = []

        with fitz.open(path) as document:
            for page_number, page in enumerate(document, start=1):
                native_text = page.get_text("text").strip()

                if native_text:
                    output.append(native_text)
                    continue

                if not self.ocr_enabled:
                    continue

                self.logger.log(
                    f"OCR required for PDF page {page_number}"
                )

                pixmap = page.get_pixmap(
                    matrix=fitz.Matrix(2, 2),
                    alpha=False,
                )

                image_bytes = pixmap.tobytes(
                    "png"
                )

                image = Image.open(
                    io.BytesIO(image_bytes)
                )

                ocr_text = pytesseract.image_to_string(
                    image
                ).strip()

                if ocr_text:
                    output.append(ocr_text)

        return "\n\n".join(output)

    # ------------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------------

    def _process_docx(self, path: Path) -> str:
        if Document is None:
            err = RuntimeError(
                "python-docx is required for DOCX files."
            )
            self.logger.exception("python-docx is required for DOCX files.", err)
                       
            raise err

        document = Document(path)

        parts: list[str] = []

        # Paragraphs.
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                parts.append(text)

        # Tables.
        for table in document.tables:
            for row in table.rows:
                cells = [
                    cell.text.strip()
                    for cell in row.cells
                ]

                cells = [
                    cell for cell in cells
                    if cell
                ]

                if cells:
                    parts.append(
                        " | ".join(cells)
                    )

        # Images embedded in DOCX are intentionally not silently ignored.
        # python-docx exposes relationships, but extracting every image
        # requires dealing with the underlying package. We use a temporary
        # extraction pass and OCR the images.
        if self.ocr_enabled:
            image_text = self._extract_docx_images(path)

            if image_text:
                parts.append(image_text)

        return "\n\n".join(parts)

    def _extract_docx_images(self, path: Path) -> str:
        """
        Extract images from the DOCX ZIP package and OCR them.

        DOCX is an OpenXML ZIP container, so this avoids adding another
        dependency merely to retrieve media files.
        """

        import zipfile

        output: list[str] = []

        with zipfile.ZipFile(path) as archive:
            image_names = [
                name
                for name in archive.namelist()
                if name.startswith("word/media/")
            ]

            for image_name in image_names:
                try:
                    image_data = archive.read(image_name)

                    image = Image.open(
                        io.BytesIO(image_data)
                    )

                    if not self._image_is_usable(image):
                        continue

                    text = pytesseract.image_to_string(
                        image
                    ).strip()

                    if text:
                        output.append(text)

                except Exception as exc:
                    self.logger.warn(
                        f"Failed to OCR embedded DOCX image "
                        f"{image_name}: {exc}"
                    )

        return "\n\n".join(output)

    # ------------------------------------------------------------------------
    # Legacy DOC
    # ------------------------------------------------------------------------

    def _process_doc(self, path: Path) -> str:
        """
        Legacy .doc files are not reliably supported by python-docx.

        First attempt native conversion using LibreOffice/soffice, then
        process the resulting DOCX.

        If conversion is unavailable, fail clearly.
        """

        converted = self._convert_with_libreoffice(
            path,
            target_format="docx",
        )

        if converted is None:
            raise RuntimeError(
                "Unable to process .doc file. "
                "Install LibreOffice/soffice for legacy DOC conversion."
            )

        try:
            return self._process_docx(converted)
        finally:
            self._safe_remove(converted)

    # ------------------------------------------------------------------------
    # PowerPoint
    # ------------------------------------------------------------------------

    def _process_pptx(self, path: Path) -> str:
        """
        Native PPTX parsing first.

        Extracts:
          - slide text
          - table contents
          - speaker notes where available
          - OCR from embedded images
        """

        if Presentation is None:
            raise RuntimeError(
                "python-pptx is required for PPTX files."
            )

        presentation = Presentation(path)

        slides: list[str] = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_parts: list[str] = [
                f"Slide {slide_number}"
            ]

            for shape in slide.shapes:
                # Normal text.
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slide_parts.append(text)

                # Tables.
                if getattr(
                    shape,
                    "has_table",
                    False,
                ):
                    for row in shape.table.rows:
                        cells = [
                            cell.text.strip()
                            for cell in row.cells
                        ]

                        if any(cells):
                            slide_parts.append(
                                " | ".join(cells)
                            )

                # Embedded images.
                if (
                    self.ocr_enabled
                    and getattr(shape, "shape_type", None) == 13
                ):
                    try:
                        image = Image.open(
                            io.BytesIO(
                                shape.image.blob
                            )
                        )

                        if self._image_is_usable(image):
                            text = pytesseract.image_to_string(
                                image
                            ).strip()

                            if text:
                                slide_parts.append(text)

                    except Exception as exc:
                        self.logger.warn(
                            f"Failed to OCR image on "
                            f"slide {slide_number}: {exc}"
                        )

            # Notes are optional in python-pptx versions.
            try:
                notes_slide = slide.notes_slide

                notes_text: list[str] = []

                for shape in notes_slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()

                        if text:
                            notes_text.append(text)

                if notes_text:
                    slide_parts.append(
                        "\n".join(notes_text)
                    )

            except Exception:
                pass

            slides.append(
                "\n".join(slide_parts)
            )

        return "\n\n".join(slides)

    def _process_ppt(self, path: Path) -> str:
        """
        Legacy .ppt:

        1. Try native/available conversion.
        2. If conversion succeeds, process resulting PPTX.
        3. Otherwise convert to PDF and use the PDF parser.
        """

        # LibreOffice can convert legacy PPT directly to PPTX.
        converted_pptx = self._convert_with_libreoffice(
            path,
            target_format="pptx",
        )

        if converted_pptx is not None:
            try:
                return self._process_pptx(
                    converted_pptx
                )
            finally:
                self._safe_remove(converted_pptx)

        # Required fallback: convert to PDF.
        converted_pdf = self._convert_with_libreoffice(
            path,
            target_format="pdf",
        )

        if converted_pdf is not None:
            try:
                return self._process_pdf(
                    converted_pdf
                )
            finally:
                self._safe_remove(converted_pdf)

        raise RuntimeError(
            "Unable to process .ppt file. "
            "Install LibreOffice/soffice for legacy PPT conversion."
        )

    # ------------------------------------------------------------------------
    # Images
    # ------------------------------------------------------------------------

    def _process_image(self, path: Path) -> str:
        """
        Standalone image OCR.
        """

        if not self.ocr_enabled:
            raise RuntimeError(
                "OCR is disabled."
            )

        with Image.open(path) as image:
            if not self._image_is_usable(image):
                return ""

            # Tesseract works more reliably with RGB/RGBA converted to RGB.
            image = image.convert("RGB")

            return pytesseract.image_to_string(
                image
            ).strip()

    @staticmethod
    def _image_is_usable(image: Image.Image) -> bool:
        width, height = image.size

        return (
            width >= MIN_IMAGE_WIDTH
            and height >= MIN_IMAGE_HEIGHT
        )

    # ------------------------------------------------------------------------
    # Plain text / Markdown / Logs
    # ------------------------------------------------------------------------

    def _process_text(self, path: Path) -> str:
        return self._read_text(path)

    def _process_markdown(self, path: Path) -> str:
        text = self._read_text(path)

        # Preserve headings and code blocks while removing only obvious
        # presentation syntax.
        text = re.sub(
            r"!\[([^\]]*)\]\([^)]+\)",
            r"\1",
            text,
        )

        text = re.sub(
            r"\[([^\]]+)\]\([^)]+\)",
            r"\1",
            text,
        )

        text = re.sub(
            r"^#{1,6}\s*",
            "",
            text,
            flags=re.MULTILINE,
        )

        return text

    # ------------------------------------------------------------------------
    # HTML
    # ------------------------------------------------------------------------

    def _process_html(self, path: Path) -> str:
        html = self._read_text(path)

        # BeautifulSoup is strongly recommended for real HTML.
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            # Reasonable fallback if dependency isn't installed.
            html = re.sub(
                r"<script\b[^>]*>.*?</script>",
                "",
                html,
                flags=re.IGNORECASE | re.DOTALL,
            )

            html = re.sub(
                r"<style\b[^>]*>.*?</style>",
                "",
                html,
                flags=re.IGNORECASE | re.DOTALL,
            )

            html = re.sub(
                r"<[^>]+>",
                " ",
                html,
            )

            return html

        soup = BeautifulSoup(
            html,
            "html.parser",
        )

        # Remove content that generally should not enter the RAG corpus.
        for element in soup(
            ["script", "style", "noscript", "svg"]
        ):
            element.decompose()

        return soup.get_text(
            separator="\n",
            strip=True,
        )

    # ------------------------------------------------------------------------
    # RTF
    # ------------------------------------------------------------------------

    def _process_rtf(self, path: Path) -> str:
        if rtf_to_text is None:
            raise RuntimeError(
                "striprtf is required for RTF files."
            )

        raw = self._read_text(
            path,
            encoding="utf-8",
            errors="ignore",
        )

        return rtf_to_text(raw)

    # ------------------------------------------------------------------------
    # LaTeX
    # ------------------------------------------------------------------------

    def _process_latex(self, path: Path) -> str:
        text = self._read_text(path)

        # Remove comments but preserve escaped \%.
        text = re.sub(
            r"(?<!\\)%.*$",
            "",
            text,
            flags=re.MULTILINE,
        )

        # Remove common environments that generally do not contain useful
        # retrieval text.
        text = re.sub(
            r"\\begin\{(?:equation|equation\*|align|align\*)\}.*?"
            r"\\end\{(?:equation|equation\*|align|align\*)\}",
            "",
            text,
            flags=re.DOTALL,
        )

        # Convert common structural commands to whitespace.
        text = re.sub(
            r"\\(?:section|subsection|subsubsection|paragraph)"
            r"\*?\{([^}]*)\}",
            r"\1\n",
            text,
        )

        # Remove commands while retaining their arguments when practical.
        text = re.sub(
            r"\\[a-zA-Z]+\*?(?:\[[^\]]*\])?",
            "",
            text,
        )

        # Remove braces used by LaTeX syntax.
        text = text.replace("{", "")
        text = text.replace("}", "")

        # Basic escaped characters.
        replacements = {
            r"\&": "&",
            r"\%": "%",
            r"\$": "$",
            r"\#": "#",
            r"\_": "_",
            r"\{": "{",
            r"\}": "}",
            r"~": " ",
        }

        for source, target in replacements.items():
            text = text.replace(source, target)

        return text

    # ------------------------------------------------------------------------
    # XML
    # ------------------------------------------------------------------------

    def _process_xml(self, path: Path) -> str:
        try:
            root = ET.parse(path).getroot()

            parts = [
                value.strip()
                for value in root.itertext()
                if value.strip()
            ]

            return "\n".join(parts)

        except ET.ParseError:
            # Some XML-ish files are malformed. Falling back to text is
            # preferable to silently dropping the document.
            self.logger.warn(
                f"Malformed XML detected in {path.name}; "
                f"falling back to raw text."
            )

            return self._read_text(path)

    # ------------------------------------------------------------------------
    # JSON
    # ------------------------------------------------------------------------

    def _process_json(self, path: Path) -> str:
        text = self._read_text(path)

        try:
            data = json.loads(text)

            return json.dumps(
                data,
                indent=2,
                ensure_ascii=False,
            )

        except json.JSONDecodeError:
            self.logger.warn(
                f"Invalid JSON in {path.name}; "
                f"falling back to raw text."
            )

            return text

    # ------------------------------------------------------------------------
    # CSV / TSV
    # ------------------------------------------------------------------------

    def _process_csv(self, path: Path) -> str:
        return self._process_delimited(
            path,
            delimiter=",",
        )

    def _process_tsv(self, path: Path) -> str:
        return self._process_delimited(
            path,
            delimiter="\t",
        )

    def _process_delimited(
        self,
        path: Path,
        delimiter: str,
    ) -> str:
        text = self._read_text(path)

        reader = csv.reader(
            io.StringIO(text),
            delimiter=delimiter,
        )

        rows: list[str] = []

        for row in reader:
            cleaned = [
                cell.strip()
                for cell in row
            ]

            if any(cleaned):
                rows.append(
                    " | ".join(cleaned)
                )

        return "\n".join(rows)

    # ------------------------------------------------------------------------
    # Excel
    # ------------------------------------------------------------------------

    def _process_xlsx(self, path: Path) -> str:
        if openpyxl is None:
            raise RuntimeError(
                "openpyxl is required for XLSX files."
            )

        workbook = openpyxl.load_workbook(
            path,
            read_only=True,
            data_only=True,
        )

        parts: list[str] = []

        try:
            for worksheet in workbook.worksheets:
                parts.append(
                    f"Sheet: {worksheet.title}"
                )

                for row in worksheet.iter_rows(
                    values_only=True
                ):
                    values = [
                        ""
                        if value is None
                        else str(value).strip()
                        for value in row
                    ]

                    if any(values):
                        parts.append(
                            " | ".join(values)
                        )

                parts.append("")

        finally:
            workbook.close()

        return "\n".join(parts)

    def _process_xls(self, path: Path) -> str:
        """
        Legacy XLS.

        Convert to XLSX using LibreOffice, then use the native XLSX parser.
        """

        converted = self._convert_with_libreoffice(
            path,
            target_format="xlsx",
        )

        if converted is None:
            raise RuntimeError(
                "Unable to process .xls file. "
                "Install LibreOffice/soffice."
            )

        try:
            return self._process_xlsx(converted)
        finally:
            self._safe_remove(converted)

    # ------------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------------

    @staticmethod
    def _read_text(
        path: Path,
        encoding: str = "utf-8",
        errors: str = "replace",
    ) -> str:
        return path.read_text(
            encoding=encoding,
            errors=errors,
        )

    def _convert_with_libreoffice(
        self,
        path: Path,
        target_format: str,
    ) -> Optional[Path]:
        """
        Convert a document using LibreOffice in headless mode.

        Returns:
            Path to converted file, or None if LibreOffice is unavailable or
            conversion failed.

        This is intentionally isolated so the rest of the parser remains
        platform-independent.
        """

        self.logger.log(f"Converting {path} to {target_format}")

        soffice = (
            shutil.which("soffice")
            or shutil.which("libreoffice")
        )

        if soffice is None:
            self.logger.warn(
                "LibreOffice/soffice not found; "
                f"cannot convert {path.name}."
            )
            return None

        temp_dir = Path(
            tempfile.mkdtemp(
                prefix="rag_conversion_"
            )
        )

        try:
            command = [
                soffice,
                "--headless",
                "--convert-to",
                target_format,
                "--outdir",
                str(temp_dir),
                str(path),
            ]

            result = subprocess.run(
                command,
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )

            if result.returncode != 0:
                self.logger.warn(
                    f"LibreOffice conversion failed for "
                    f"{path.name}: {result.stderr.strip()}"
                )

                shutil.rmtree(
                    temp_dir,
                    ignore_errors=True,
                )

                return None

            expected = (
                temp_dir
                / f"{path.stem}.{target_format}"
            )

            if not expected.exists():
                # LibreOffice may normalise filenames differently.
                candidates = list(
                    temp_dir.glob(
                        f"*.{target_format}"
                    )
                )

                if not candidates:
                    shutil.rmtree(
                        temp_dir,
                        ignore_errors=True,
                    )

                    return None

                expected = candidates[0]

            return expected

        except (
            subprocess.TimeoutExpired,
            OSError,
        ) as exc:
            self.logger.warn(
                f"Conversion failed for {path.name}: {exc}"
            )

            shutil.rmtree(
                temp_dir,
                ignore_errors=True,
            )

            return None

    @staticmethod
    def _safe_remove(path: Optional[Path]) -> None:
        if path is None:
            return

        try:
            if path.is_dir():
                shutil.rmtree(
                    path,
                    ignore_errors=True,
                )
            elif path.exists():
                path.unlink()
        except OSError:
            pass


# ============================================================================
# Simple functional API
# ============================================================================

def extract_text(
    file_path: str | Path,
    logger=None,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> list[str]:
    """
    Backwards-friendly convenience function.

    Returns only chunks, exactly like the original API.
    """

    parser = DocumentParser(
        logger=logger,
        chunk_size=chunk_size,
        overlap=overlap,
    )

    return parser.process(file_path)
